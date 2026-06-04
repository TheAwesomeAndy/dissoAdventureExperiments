#!/usr/bin/env python3
"""
Build the ARSPI-Net PhD defense PowerPoint and the matching word-for-word
spoken transcript from a single source of truth (the SLIDES list below).

- Repo result figures (PDF) are rasterized on demand (PyMuPDF), cached in figpng/.
- Equations / tables / TikZ diagrams come from latexgen.py (figpng/*.png).
- External images live in assets/.
- Every slide carries speaker notes; the same narration is emitted to
  defense/spoken_transcript.md.

Run:  python3 defense_build/latexgen.py  &&  python3 defense_build/build_deck.py
"""
import os, hashlib
from PIL import Image
import fitz
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
REPO = os.path.dirname(HERE)
FIGPNG = os.path.join(HERE, "figpng")
ASSETS = os.path.join(HERE, "assets")
OUT_PPTX = os.environ.get("DECK_OUT", os.path.join(REPO, "defense", "ARSPINet_Defense.pptx"))
OUT_TXT = os.path.join(REPO, "defense", "spoken_transcript.md")
NO_MOVIE = bool(os.environ.get("NO_MOVIE"))  # verification builds use static posters

# ---- Stony Brook palette -----------------------------------------------------
NAVY = RGBColor(0x1A, 0x2A, 0x4F)
TEAL = RGBColor(0x1C, 0x72, 0x93)
ACCENT = RGBColor(0xB5, 0x65, 0x1D)
GREY = RGBColor(0x8A, 0x94, 0xA6)
DARK = RGBColor(0x22, 0x2A, 0x36)
LIGHT = RGBColor(0xEA, 0xF1, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Calibri"
FONT_H = "Calibri"

SW, SH = 13.333, 7.5
missing = []

# ---- image helpers -----------------------------------------------------------
def pdf_to_png(relpath, dpi=300):
    src = relpath if os.path.isabs(relpath) else os.path.join(REPO, relpath)
    if not os.path.exists(src):
        missing.append(relpath); return None
    h = hashlib.md5((relpath + str(dpi)).encode()).hexdigest()[:10]
    out = os.path.join(FIGPNG, "repo_" + h + ".png")
    if not os.path.exists(out):
        d = fitz.open(src)
        pix = d[0].get_pixmap(matrix=fitz.Matrix(dpi/72, dpi/72))
        pix.save(out); d.close()
    return out

def resolve(token):
    """Map a content token to an actual image path."""
    if token is None:
        return None
    if token.endswith(".pdf"):
        return pdf_to_png(token)
    if "/" in token:                      # explicit repo-relative raster path
        p = token if os.path.isabs(token) else os.path.join(REPO, token)
        if os.path.exists(p):
            return p
        missing.append(token); return None
    # bare name -> rendered latex/tikz/table asset (figpng) or external asset
    for cand in (os.path.join(FIGPNG, token), os.path.join(ASSETS, token),
                 os.path.join(FIGPNG, token + ".png"),
                 os.path.join(ASSETS, token + ".png"),
                 os.path.join(ASSETS, token + ".jpg")):
        if os.path.exists(cand):
            return cand
    missing.append(token); return None

def imsize(path):
    with Image.open(path) as im:
        return im.size

def fit(path, l, t, w, h):
    iw, ih = imsize(path)
    ar, bar = iw/ih, w/h
    if ar > bar:
        nw, nh = w, w/ar
    else:
        nh, nw = h, h*ar
    return Inches(l + (w-nw)/2), Inches(t + (h-nh)/2), Inches(nw), Inches(nh)

# ---- low-level slide primitives ---------------------------------------------
def textbox(slide, text, l, t, w, h, size=18, color=DARK, bold=False, italic=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, line_spacing=1.0):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.name = font; r.font.color.rgb = color
    return tb

def rect(slide, l, t, w, h, fill):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill; sp.line.fill.background()
    sp.shadow.inherit = False
    return sp

def bullets(slide, items, l, t, w, h, size=19, color=DARK, gap=8):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.05
        r = p.add_run(); r.text = "▸  " + it
        r.font.size = Pt(size); r.font.name = FONT; r.font.color.rgb = color
    return tb

def picture(slide, path, l, t, w, h, caption=None, credit=None):
    if not path:
        textbox(slide, "[missing figure]", l, t, w, h, 14, GREY, italic=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return
    cap_h = 0.3 if caption else 0.0
    L, T, W, H = fit(path, l, t, w, h - cap_h)
    slide.shapes.add_picture(path, L, T, W, H)
    if caption:
        textbox(slide, caption, l, t + h - cap_h, w, cap_h, 11.5, GREY, italic=True,
                align=PP_ALIGN.CENTER)
    if credit:
        textbox(slide, credit, l, t + h - 0.22, w, 0.2, 8.5, GREY, align=PP_ALIGN.RIGHT)

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# ---- slide templates ---------------------------------------------------------
def s_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def titlebar(slide, title):
    rect(slide, 0, 0, SW, 1.02, NAVY)
    rect(slide, 0, 1.02, SW, 0.06, ACCENT)
    textbox(slide, title, 0.55, 0, SW-1.1, 1.02, 26, WHITE, bold=True,
            anchor=MSO_ANCHOR.MIDDLE, font=FONT_H)

def footer(slide, n):
    textbox(slide, "ARSPI-Net   ·   A. Lane", 0.5, 7.16, 5, 0.3, 9.5, GREY)
    textbox(slide, str(n), SW-1.1, 7.16, 0.6, 0.3, 9.5, GREY, align=PP_ALIGN.RIGHT)

def render(prs, sl, n):
    k = sl["kind"]
    slide = s_blank(prs)
    if k == "title":
        rect(slide, 0, 0, SW, SH, NAVY)
        rect(slide, 0, 3.05, SW, 0.05, ACCENT)
        textbox(slide, sl["title"], 0.8, 1.5, SW-1.6, 1.3, 60, WHITE, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_H)
        textbox(slide, sl["sub"], 0.8, 3.2, SW-1.6, 1.4, 23, LIGHT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        textbox(slide, sl["foot"], 0.8, 5.7, SW-1.6, 1.2, 18, GREY,
                align=PP_ALIGN.CENTER)
    elif k == "section":
        rect(slide, 0, 0, SW, SH, NAVY)
        if sl.get("kicker"):
            textbox(slide, sl["kicker"], 0.9, 2.7, SW-1.8, 0.6, 18, ACCENT, bold=True,
                    align=PP_ALIGN.CENTER)
        textbox(slide, sl["title"], 0.9, 3.2, SW-1.8, 1.4, 40, WHITE, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_H)
    elif k == "statement":
        rect(slide, 0, 0, 0.28, SH, ACCENT)
        if sl.get("kicker"):
            textbox(slide, sl["kicker"], 1.0, 1.5, SW-2.0, 0.6, 18, TEAL, bold=True)
        textbox(slide, sl["title"], 1.0, 2.0, SW-2.0, 3.2, 36, NAVY, bold=True,
                anchor=MSO_ANCHOR.MIDDLE, font=FONT_H, line_spacing=1.05)
        if sl.get("sub"):
            textbox(slide, sl["sub"], 1.0, 5.2, SW-2.0, 1.3, 20, DARK, italic=True)
        footer(slide, n)
    elif k == "close":
        rect(slide, 0, 0, SW, SH, NAVY)
        textbox(slide, sl["title"], 0.8, 2.4, SW-1.6, 1.6, 48, WHITE, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_H)
        textbox(slide, sl["sub"], 0.8, 4.2, SW-1.6, 1.6, 20, LIGHT, align=PP_ALIGN.CENTER)
    else:
        # content slide: title bar + optional headline + visuals/text
        titlebar(slide, sl["title"])
        top = 1.32
        if sl.get("headline"):
            textbox(slide, sl["headline"], 0.55, top, SW-1.1, 0.8, 21, ACCENT, bold=True,
                    anchor=MSO_ANCHOR.MIDDLE)
            top += 0.92
        bottom = 7.05
        area = (0.55, top, SW-1.1, bottom-top)
        if k == "figure":
            picture(slide, resolve(sl["img"]), *area, caption=sl.get("caption"),
                    credit=sl.get("credit"))
        elif k == "figtext":
            fw = (SW-1.1) * 0.58
            picture(slide, resolve(sl["img"]), 0.55, top, fw, bottom-top,
                    caption=sl.get("caption"), credit=sl.get("credit"))
            bullets(slide, sl["bullets"], 0.55+fw+0.3, top, (SW-1.1)-fw-0.3, bottom-top,
                    size=sl.get("bsize", 19))
        elif k == "eqfig":
            eqh = 1.25
            picture(slide, resolve(sl["img"]), 0.55, top, SW-1.1, (bottom-top)-eqh,
                    caption=sl.get("caption"), credit=sl.get("credit"))
            picture(slide, resolve(sl["eq"]), 1.8, bottom-eqh, SW-3.6, eqh-0.05)
        elif k == "eq":
            eqh = 1.7
            picture(slide, resolve(sl["eq"]), 1.4, top+0.1, SW-2.8, eqh)
            if sl.get("bullets"):
                bullets(slide, sl["bullets"], 1.2, top+eqh+0.25, SW-2.4, bottom-(top+eqh+0.25),
                        size=sl.get("bsize", 19))
        elif k == "twofig":
            gap = 0.4; cw = (SW-1.1-gap)/2
            picture(slide, resolve(sl["img"]), 0.55, top, cw, bottom-top,
                    caption=sl.get("caption"), credit=sl.get("credit"))
            picture(slide, resolve(sl["img2"]), 0.55+cw+gap, top, cw, bottom-top,
                    caption=sl.get("caption2"), credit=sl.get("credit2"))
        elif k == "table":
            picture(slide, resolve(sl["img"]), 0.9, top, SW-1.8, bottom-top,
                    caption=sl.get("caption"))
        elif k == "tabletext":
            tw = (SW-1.1) * 0.6
            picture(slide, resolve(sl["img"]), 0.55, top, tw, bottom-top, caption=sl.get("caption"))
            bullets(slide, sl["bullets"], 0.55+tw+0.3, top, (SW-1.1)-tw-0.3, bottom-top,
                    size=sl.get("bsize", 18))
        elif k == "bullets":
            bullets(slide, sl["bullets"], 0.8, top, SW-1.6, bottom-top, size=sl.get("bsize", 22))
        elif k == "movie":
            poster = pdf_to_png(sl["poster"]) if sl["poster"].endswith(".pdf") else resolve(sl["poster"])
            mv = os.path.join(REPO, sl["movie"])
            if os.path.exists(mv) and not NO_MOVIE:
                L, T, W, H = fit(poster, 0.55, top, (SW-1.1)*0.62, bottom-top)
                slide.shapes.add_movie(mv, L, T, W, H, poster_frame_image=poster, mime_type="video/mp4")
            else:
                picture(slide, poster, 0.55, top, (SW-1.1)*0.62, bottom-top)
            if sl.get("caption"):
                textbox(slide, sl["caption"], 0.55, bottom-0.32, (SW-1.1)*0.62, 0.3, 11.5, GREY,
                        italic=True, align=PP_ALIGN.CENTER)
            if sl.get("bullets"):
                bx = 0.55 + (SW-1.1)*0.62 + 0.3
                bullets(slide, sl["bullets"], bx, top, (SW-0.55)-bx, bottom-top, size=18)
        footer(slide, n)
    notes(slide, sl.get("say", ""))
    return slide

# =============================================================================
#  CONTENT  (every "say" is the word-for-word spoken script — human, confident)
# =============================================================================
SLIDES = [
 dict(kind="title", title="ARSPI-Net",
   sub="Affective Reservoir-Spike Processing and Inference Network",
   foot="A four-level interpretable neuromorphic framework for clinical EEG analysis\n\nAndrew Lane     ·     Electrical & Computer Engineering, Stony Brook University     ·     Dissertation Defense · 2026",
   say="Good morning, and thank you all for being here. I'm not going to open with an architecture. I'm going to open with a decision a doctor has to make, because that decision is the reason this whole project exists."),

 dict(kind="statement", kicker="THE DECISION",
   title="An AI reads the scan. It's 89% sure of the diagnosis.\nIt even names the drug. Then you ask why — and it has nothing.",
   sub="A prediction you can't interrogate is not something a clinician can act on.",
   say="Picture a psychiatrist with a patient who has severe, treatment-resistant depression. An AI tool reads the patient's EEG and returns a diagnosis, eighty-nine percent confident, and it even recommends a specific drug. That sounds like the future of medicine — until the doctor asks the one question that matters. Why? What in this brain led to that number? And the model has nothing to say, because it doesn't know. It only knows the pattern matched."),

 dict(kind="figure", title="The same machine, the same confidence", img="goldfinch.jpg",
   headline="“American goldfinch, 87%” — and it's just as sure about something that isn't a bird.",
   caption="American goldfinch", credit="Photo: Wikimedia Commons, CC BY-SA 4.0",
   say="Here's what I mean. This is the same kind of model that looks at this photo and says American goldfinch, eighty-seven percent — and it's right. But show it something that plainly isn't a bird, and it is just as confident. From the outside, right for the right reasons and right for the wrong reasons look identical. In a clinic, that isn't a curiosity. It's a liability."),

 dict(kind="statement", kicker="THE QUESTION",
   title="Can we recover the brain's hidden dynamics from a\nnoisy, low-dimensional trace — and keep an explanation\na clinician can actually read?",
   sub="That is the whole talk. Everything after this is how I built a model that answers yes to both halves.",
   say="So here is the question I set out to answer. Can we recover the hidden dynamics of the brain from a noisy, low-dimensional trace — and do it with a model whose internal workings a clinician can actually read? Both halves matter. Accuracy alone isn't enough, and an explanation of a wrong mechanism is worse than none. Everything from here is how I built a system that answers yes to both."),

 dict(kind="section", kicker="PART ONE", title="The measurement problem",
   say="So let me start with the signal itself — where it comes from, and why a clean-looking trace is so deceptively hard to read."),

 dict(kind="figure", title="Recovering a hidden system from one trace", img="insight_mars.jpg",
   headline="NASA's InSight read a planet's deep interior from a single seismometer on the surface.",
   credit="Image: NASA/JPL-Caltech, public domain",
   say="Recovering a hidden system from a thin signal is not a fantasy — physics does it all the time. NASA put a single seismometer on the surface of Mars, and from that one squiggle reconstructed the planet's deep interior, structure no instrument could ever touch directly. The trick is that they knew the physics connecting the inside to the surface. That is exactly the bet I'm making with the brain."),

 dict(kind="eqfig", title="Where the EEG signal comes from", img="p300_erp.png",
   caption="What reaches the scalp: averaged event-related potentials", credit="ERP waveforms: Wikimedia Commons, CC BY-SA",
   eq="eq_observation",
   say="An EEG signal is born when thousands of cortical pyramidal cells fire together and act like a tiny electrical dipole. What we actually record on the scalp isn't that source — it's a noisy, smeared, averaged projection of it, like these evoked waveforms. Formally: the measurement x is some mixing function g of the latent cortical response r, plus noise. The whole game is getting back to r."),

 dict(kind="eqfig", title="The volume conduction problem", img="volume_conduction.png",
   eq="eq_inverse",
   headline="Volume conduction: tissue between source and sensor spreads each dipole across many electrodes — a spatial low-pass filter.",
   caption="Pyramidal-cell dipoles → volume & capacitive conduction to the scalp electrode",
   credit="Jackson & Bolger, Psychophysiology (2014)",
   say="The reason this is hard has a name — volume conduction. A cortical source, those pyramidal cells at the bottom of the figure firing together as a dipole, doesn't reach the electrode directly. The current has to spread through brain tissue, cerebrospinal fluid, the skull, and the scalp, and every layer smears it. That is volume conduction: the head acts as a spatial low-pass filter, so a single source blurs across many electrodes and every electrode records a mixture from across the whole cortex. Recovering the source from that blurred mixture — with thirty-four channels standing in for billions of neurons — is a deeply underdetermined inverse problem. Reading those shadows is the whole challenge."),

 dict(kind="figure", title="Why EEG and not fMRI", img="eeg_fmri_resolution.png",
   headline="Affect happens in milliseconds. EEG keeps the timing; fMRI trades it for spatial detail.",
   caption="Spatial vs. temporal resolution of brain-imaging methods",
   credit="Wikimedia Commons, CC BY 4.0",
   say="People ask why not just use fMRI, with its beautiful spatial maps. Look at where the methods sit: fMRI is over on the high-spatial side, but low on the temporal axis, because it measures blood flow — a sluggish metabolic echo that lags the actual neural event by seconds. Surface EEG is the opposite corner: modest spatial resolution, but millisecond timing. And emotion and threat detection happen in milliseconds. Using fMRI for that is like following a fast piano piece by measuring the temperature of the keys — you learn which were pressed, but the music is gone. EEG keeps the timing. It gives us the music, and the discriminative information lives in that timing."),

 dict(kind="figtext", title="The data: real patients, not textbook cases", img="comorbidity.jpg",
   caption="Transdiagnostic symptom overlap", credit="Figure: PMC (HiTOP), CC BY",
   bullets=["SHAPE cohort: 211 adults, 34-channel ERPs",
            "Transdiagnostic — about three diagnoses per person",
            "MDD, PTSD, substance use, anxiety, all overlapping",
            "Group averages show the signal — and bury the individual"],
   say="And the data are real. The SHAPE cohort is two hundred eleven adults, and it's transdiagnostic — these aren't clean textbook cases. The average patient carries about three overlapping diagnoses: depression, PTSD, substance use, anxiety. We've known for decades that these conditions show up in EEG at the group level. But you only get that clean group signal by averaging away the individual — and you can't average away the person sitting in front of you. That gap is the whole clinical problem."),

 dict(kind="section", kicker="PART TWO", title="Why a black box is the wrong tool here",
   say="That sets up the real obstacle. It isn't getting an answer out of the data — modern networks do that easily. It's getting an answer you can trust."),

 dict(kind="figure", title="Post-hoc explanations explain the shortcut, not the science",
   img="cleverhans_shap.png",
   headline="A saliency map will faithfully highlight the artifact a model cheated on.",
   credit="Figure: Frontiers in AI (2025), CC BY",
   say="The standard rebuttal is that we've solved the black box — just run SHAP or a saliency map. But those explain the model from the outside; they approximate where it looked, not why. If a network secretly learned to key off the hum of the EEG machine in room four, a saliency map will faithfully, confidently highlight that artifact. It's the Clever Hans effect with a colorful heatmap on top. It can make a broken model look trustworthy, which in medicine is the most dangerous outcome of all."),

 dict(kind="figure", title="Intrinsic interpretability instead", img="pictures/chGraphNeuralNetworks/fig_shap_comparison.pdf",
   headline="EEGNet's attention peaks at 402–691 ms; ARSPI-Net lands on the named ERP windows.",
   caption="Saliency timing: black box vs. ARSPI-Net",
   say="I want the opposite: a model whose internal variables are the explanation. Here's the contrast on the same data. The black box's attention peaks late, between four hundred and seven hundred milliseconds — outside the windows that decades of clinical work tie to emotional processing. It's accurate, but for reasons a neurologist can't use. ARSPI-Net's representation sits right on the named event-related components — the P300, the early late-positive potential. One is a gradient about a model. The other is an account in the brain's own units."),

 dict(kind="table", title="The claim is the count, not the accuracy", img="tbl_landscape",
   headline="First architecture interpretable at all four levels in one pipeline.",
   say="So here's where this work sits. A black box with SHAP gives you zero intrinsic levels of interpretability. Prototype methods give one. NeuCube, the closest neighbor, gives two. ARSPI-Net exposes four, in a single pipeline. That count — not a leaderboard number — is the contribution, and I'll show you all four levels measured."),

 dict(kind="section", kicker="PART THREE", title="The approach: a spiking reservoir",
   say="So here's the machine I built to thread that needle — and why every piece of it is chosen from theory, not tuned for a score."),

 dict(kind="eqfig", title="Spikes are events in time", img="lif_neuron.jpg", eq="eq_lif",
   caption="Leaky integrate-and-fire neuron", credit="Schematic: Wikimedia Commons, CC BY-SA",
   headline="The third generation of neural computation — and it matches the physics of EEG.",
   say="The brain doesn't compute in clocked, dense arithmetic. It computes in spikes — sparse events whose timing carries the information. That's the third generation of neural models, and it's a natural match for a signal whose content is in its timing. My unit is the leaky integrate-and-fire neuron: it integrates input until it crosses threshold, fires, and resets. One parameter, the leak beta, sets how long it remembers."),

 dict(kind="figure", title="Spikes are also almost free", img="tikz_energy",
   headline="Event-driven neuromorphic silicon: ≈45× less energy per operation than dense CMOS.",
   say="There's a second reason to care about spikes, beyond the physics — energy. A dense multiply-accumulate on conventional hardware costs around nine hundred picojoules; on event-driven neuromorphic silicon the same operation is roughly twenty. That's about forty-five times less, and it's what could eventually take this off a server and onto a wearable at the bedside. Everything here ran in software, so I treat that number as motivation and a future on-chip measurement — not a result I'm claiming today."),

 dict(kind="figure", title="The reservoir: a calibrated ruler, not an oracle", img="pictures/ch2/LSMArchitecture.png",
   headline="Random recurrent weights, fixed forever. Train them and you turn the ruler into a black box.",
   caption="Liquid state machine: a fixed recurrent reservoir feeding trained linear readouts",
   credit="LSM schematic (after Maass et al.)",
   say="I wire two hundred fifty-six of these neurons into a recurrent pool — a reservoir. Here's the intuition: throw a rock into a still pond. You can't measure the rock in flight, but the ripples it leaves are a faithful record of its size and speed. The input is the rock, the reservoir is the water, and a simple readout reads the ripples. Now the decision that makes everything else in this talk possible — the recurrent weights are random, and they are fixed. I never train them. The instant you train the water, it starts bending its own physics to chase a score, and the ripples stop being an honest record of the rock. A fixed reservoir stays a calibrated instrument. And because nothing inside it is trained, its entire state is a pure, reproducible function of the input — which is exactly what lets me read it later, neuron by neuron."),

 dict(kind="eqfig", title="Why this is recoverable, not a hope", img="kernel_trick.png", eq="eq_koopman",
   headline="Cover lifts the data to separability; Takens says the geometry survives; Koopman makes the readout linear.",
   caption="Cover's theorem: a nonlinear boundary becomes linear in a higher-dimensional lift",
   credit="Kernel machine — Wikimedia Commons, CC BY-SA",
   say="And this is provably the right move, not a lucky one — which matters, because a committee will ask why this should work at all. Cover's theorem says that if you lift tangled data into a high enough dimensional space, it almost always becomes linearly separable; those are the ripples spreading out. Takens tells me the lift preserves the geometry of the underlying dynamics — and rather than just cite it, I measured the embedding dimension and found the trajectory needs fewer than sixty-four dimensions to live in. Koopman closes the loop: in the right space of observables, even a nonlinear system evolves linearly, which is the deep reason a plain linear readout is enough. Three results, one conclusion — a training-free core isn't a compromise. It's sufficient by construction."),

 dict(kind="eqfig", title="ARSPI-Net, end to end", img="pictures/chGraphNeuralNetworks/fig_paper_pipeline_overview.pdf",
   eq="eq_chain",
   say="Putting it together: the scalp signal drives the fixed reservoir; the spikes get a compact temporal code; that feeds a graph stage over the electrodes; and a linear readout produces the affective state. Five stages, and under each one sits a different, measurable level of interpretability. Nothing in the core is trained — the only fitted object in the whole system is that final linear readout."),

 dict(kind="bullets", title="Four contributions", bsize=21,
   bullets=["A neuromorphic temporal operator — a fixed spiking reservoir with a certified echo-state property",
            "A traceable spike-to-embedding code (BSC₆) aligned to clinical ERP components",
            "A design rule for graphs — message passing destroys contrast on dense electrode arrays",
            "A geometric correction — subject identity dominates; centering recovers the signal for every model"],
   say="That gives four contributions. A neuromorphic operator I can analyze as a dynamical system. A spike code that stays tied to clinical components. A design rule that says when graph message-passing helps and when it hurts. And the one I'm most excited about — a geometric correction that lifts every model in the field, not just mine. Let me take you through the evidence."),

 dict(kind="section", kicker="PART FOUR", title="What the model measures",
   say="Now the evidence. Four levels of it — and I want each one to be a number you can check, not an adjective."),

 dict(kind="figtext", title="The setup, briefly", img="pictures/chGraphNeuralNetworks/obs18_subject_gallery.pdf",
   caption="Per-subject embeddings — notice how different people look",
   bullets=["Subject-disjoint cross-validation, always",
            "Linear readouts, so differences reflect the representation",
            "Permutation- and bootstrap-tested",
            "Stable to the random seed within 1%"],
   say="A word on discipline, because it's what makes the numbers mean something. Every split is subject-disjoint — no patient appears in both train and test. I compare representations with linear readouts, so a difference reflects the representation and not a fancier classifier. Everything is permutation-tested, and the whole pipeline is stable to the random seed within one percent. Look at this gallery for a second — every subject's signature looks completely different. Hold that thought, because it comes back as the central result."),

 dict(kind="eqfig", title="The reservoir is stable — measured, not assumed",
   img="pictures/chDynamics/analysis6_1f_lyapunov.pdf", eq="eq_lyap",
   headline="Driven by 3,000+ real ERPs, the Lyapunov exponent is −0.054, negative every time.",
   say="First, stability — and I'm careful here, because this is where people wave their hands. The usual move is to say the spectral radius is below one, so the system is stable. But that's a property of the weights sitting in a drawer; it says nothing about what a real, structured signal does once it's driving the network. So I measured the thing that actually matters: the largest Lyapunov exponent under genuine ERP drive. Across more than three thousand real trials it's minus zero point zero five four, and negative on every single one. Two nearby states, given the same input, converge — the reservoir forgets where it started. That's the echo-state property as a measurement, not an assumption, and it's what licenses everything downstream."),

 dict(kind="movie", title="A system you can fully watch",
   movie="pictures/animations/lsm_membrane_raster.mp4", poster="pictures/animations/lsm_membrane_raster_poster.pdf",
   caption="256 membrane potentials responding to an ERP",
   bullets=["Nothing is trained, so the state is fully observable",
            "Input-locked response, then graded relaxation",
            "Fading memory you can see, not infer"],
   say="And because nothing is trained, I can watch the entire internal state — all two hundred fifty-six membrane potentials as a real ERP drives the pool. You see the input-locked burst, then the graded relaxation: fading memory, as a measured fact rather than an assumption. This complete observability is the substrate for every interpretability claim that follows."),

 dict(kind="eqfig", title="Level 1 — a temporal code clinicians can read",
   img="pictures/chLSMEmbeddings/coding_scheme_accuracy_comparison.pdf", eq="eq_bsc6",
   headline="Six bins land on N1, P200, P300, and the LPP — the clinician's own vocabulary.",
   say="Now the four levels. Level one is temporal traceability. The reservoir spits out a dense spike matrix, and I funnel it through a binned spike count — six windows. Six isn't arbitrary: those windows line up with the canonical ERP components clinicians already read, the N1, P200, P300, and late-positive potential. So when the model flags an anomaly, a clinician can see which window it lives in and check it against forty years of literature."),

 dict(kind="twofig", title="Every step has a name", img="pictures/chLSMEmbeddings/obs03_raw_bsc6_features.pdf",
   img2="pictures/chLSMEmbeddings/obs04_raw_embedding_space.pdf",
   caption="BSC₆ temporal features", caption2="PCA-64 embedding",
   say="From there it's fully traceable: raw channel, spikes, the six-bin code, and a sixty-four-dimensional embedding. Every arrow is a named object, and the whole chain is stable to the seed. There's no hidden layer where the meaning disappears."),

 dict(kind="eqfig", title="Level 2 — the discovery that reorganized the project",
   img="pictures/chGraphNeuralNetworks/exp01_variance_decomposition.pdf", eq="eq_rho",
   headline="Emotion is 8.7% of the variance. Identity is 62.6%. It pulls 7.2× harder.",
   say="Level two is where the whole project turned, and I almost misread it as a failure. You'd hope the emotional condition is what drives the variance in the embedding. It accounts for under nine percent. Subject identity — your cortical folding, the thickness of your skull, your baseline rhythm — accounts for sixty-three. Identity pulls more than seven times harder than the signal I'm actually after. And worse, the two aren't in separate directions I could just project away; they share the same principal axes, so any filter that removes identity takes the emotion with it. For a while I thought the data was telling me there was nothing there. It was telling me the opposite — the signal was real, just hiding underneath a much larger one."),

 dict(kind="eq", title="The fix is geometric, and it's label-free", eq="eq_centering",
   bullets=["Compute each subject's own mean in the 64-D space",
            "Subtract it — shift their origin to the center",
            "Identity offset gone; the affective direction untouched",
            "No labels used — this is geometry, not a classifier"],
   say="The fix is almost embarrassingly simple. For each patient, I compute their own mean in the embedding space and subtract it — I move their origin to the center. That erases the identity offset while leaving the direction of the emotional response completely intact. It uses no labels. It's pure geometry. And seeing it was only possible because the space was transparent enough to look at."),

 dict(kind="table", title="Centering is the dominant move — for everyone", img="tbl_baseline",
   headline="+13 to +21 points for every model. EEGNet 72 → 89. ARSPI-Net 59 → 79.",
   say="And here's the part that still gets me. The fix — re-centering each subject on their own mean — doesn't just help my model. Read down the centered column: every single model jumps by thirteen to twenty-one points. EEGNet goes from seventy-two to eighty-nine; mine from fifty-nine to seventy-nine. The emotional signal was in the data the whole time, in every representation, just buried under subject geometry. After centering, my training-free reservoir even ties a trained GRU with zero trainable recurrent parameters. A transparent model is what let me see that geometry and name the problem — and naming it raised the ceiling for the entire field, black boxes included. The geometric correction matters more than the architecture race — and that's the case for interpretability in a single number."),

 dict(kind="figure", title="Why it works", img="pictures/chGraphNeuralNetworks/obs14_within_between_geometry.pdf",
   headline="Before centering, same-condition trials sit farther apart than same-subject trials. Centering inverts it.",
   say="Mechanically: before centering, two trials from the same person under different emotions are closer together than two trials of the same emotion from different people. The ratio is above one — identity wins. Subtract each subject's mean and it flips below one — now the conditions separate. Three-class accuracy moves from sixty-three to seventy-nine on that one geometric move."),

 dict(kind="figure", title="A boundary that tells you something", img="pictures/appendixA/ch5_4class_fig01_gnn_comparison.pdf",
   headline="Reservoir wins by +12.5 points on the 3-class task, loses 6 on the 4-class. That's the memory–nonlinearity trade, located.",
   say="One more finding I didn't expect. When I push from three classes to four, the ranking flips — the reservoir's advantage of twelve and a half points becomes a six-point deficit. That's not a bug; it's the memory-versus-nonlinearity trade-off showing me exactly where this representation pays off: when the structure lives in timing, not in amplitude. The boundary itself is information."),

 dict(kind="figure", title="Level 3 — seven named observables",
   img="pictures/chDynamics/fig01_condition_effects.pdf",
   headline="All seven distinguish emotional from neutral input — and they're physical quantities, not learned features.",
   say="Level three reads the reservoir's dynamics directly, through seven named descriptors — firing rate, temporal sparsity, a memory timescale, and so on. These are physical quantities, not learned features, and that distinction matters: each one is something a physiologist can actually argue with. All seven separate emotional from neutral input at the subject level. And one of them, the autocorrelation timescale, has a direct clinical reading — an abnormally slow decay is the network failing to downregulate after a stressor, which is exactly the dynamics you'd expect from rumination. So the descriptor isn't just predictive; it's a hypothesis about mechanism."),

 dict(kind="twofig", title="Timing carries the signal", img="pictures/chDynamics/fig02_metric_families.pdf",
   img2="pictures/chDynamics/fig_latent_axis.pdf",
   caption="Temporal family ≈ 2.4× the amplitude family", caption2="One excitability–persistence axis",
   say="And when I group those descriptors, the temporal-structure family carries about two and a half times the signal of the amplitude family — independent confirmation that timing is where the information is. Better still, the seven collapse onto a single axis, from excitable to persistent. So the model doesn't just stamp a label on a patient; it locates them on a continuous, interpretable axis. That's the instrument, made concrete."),

 dict(kind="eqfig", title="Level 4 — message passing fails here, and the theory said it would",
   img="pictures/chGraphNeuralNetworks/exp03_propagation_operating_characteristic.pdf", eq="eq_prop",
   headline="On a small, dense electrode graph, every propagation operator loses to no propagation.",
   say="Level four is the graph across electrodes — and I expected message passing, the field's default, to help. But that operator is mathematically a low-pass filter, and over-smoothing theory makes a sharp prediction for a small, dense montage like ours: the contrast between channels should collapse within a couple of steps. So I swept every operator I could — plain smoothing, residual connections, attention — and every one of them loses to doing no propagation at all. I went looking for a tool and found a wall. But the wall is the theory being confirmed, and it turns into a design rule the EEG-graph literature didn't actually have."),

 dict(kind="figure", title="The mechanism, measured", img="pictures/chGraphNeuralNetworks/exp03b_diffusion_overlay.pdf",
   headline="Dirichlet energy drops 84% by depth two — exactly where accuracy falls fastest.",
   say="And I can show the mechanism on the real features: as you stack graph layers, the Dirichlet energy — the contrast between channels — drops eighty-four percent by depth two, and the channels become nearly identical. That's the same depth where accuracy falls off a cliff. The diagnosis and the symptom line up."),

 dict(kind="eqfig", title="So is the system actually one coherent thing?",
   img="pictures/chSynthesis/fig7_A2_kappa_vs_null.pdf", eq="eq_kappa",
   headline="Local dynamics and global topology are coupled: κ = 0.273, p < 0.001 across 211 subjects.",
   say="So instead of forcing the graph to classify, I use it to ask a deeper question: are the local dynamics and the global topology two views of one system, or two unrelated things? I measure their coupling — kappa — against a null where I shuffle the electrode labels. The answer is that they're genuinely linked: a median of zero point two seven three, significant across all two hundred eleven subjects. It's a consistent system-wide effect rather than a strong local one, and I report it exactly that way — but it means the pieces of this pipeline are describing one coherent object, which is the last thing I'd otherwise be able to claim."),

 dict(kind="table", title="Four levels, all measured", img="tbl_taxonomy",
   headline="One pipeline, interpretable at every stage.",
   say="So there are the four levels, each one a number, not an adjective. Temporal traceability correlates with the ERP at point eight two. Geometric transparency is the seven-point-two ratio and the centering lift. The dynamical descriptors track the ERP up to point eight four. And the systems coupling clears its null at high significance. To my knowledge this is the first neuromorphic EEG model to put a measurement on all four in a single pipeline."),

 dict(kind="tabletext", title="Different disorders live in different layers",
   img="tbl_disorder",
   bullets=["Substance use: dynamical descriptors, p = 0.0004 — survives correction",
            "The others point in suggestive directions",
            "We're testing them now with the SHAPE clinical team",
            "A single score can't even ask this question"],
   say="And because the model decomposes into named layers, I can ask which layer is most sensitive to which condition. Substance use separates most strongly in the dynamical descriptors, at p equals zero point zero zero zero four, and it survives correction — that's the solid clinical result. The others point in suggestive directions we're now testing with our clinical collaborators. The point isn't any single p-value; it's that a black-box score can't even pose this question. This model can."),

 dict(kind="section", kicker="PART FIVE", title="What it means, and where it goes",
   say="Let me pull all of this together and say plainly what it adds up to — and what it doesn't."),

 dict(kind="statement", kicker="THE TRADE",
   title="78.8% with a full decomposition,\nor 89.1% from a box you can't open.",
   sub="Matching a trained GRU with zero trainable recurrent parameters — and a reason for every number. That trade is the contribution.",
   say="Let me be direct about the trade at the center of this, because it's the obvious line of attack. My model gets seventy-nine percent with a complete, four-level account of why. The black box gets eighty-nine with nothing you can open. And I'll defend the seventy-nine — because in a clinic, an eighty-nine that can't justify itself is exactly the liability we opened with, while a seventy-nine you can trace end to end is an instrument a doctor can stand behind. It isn't even a large gap, and my training-free reservoir ties a fully trained recurrent network with no trained recurrence at all. Those points are the price of a glass box, and for this problem it's a price worth paying."),

 dict(kind="figure", title="The biological prism", img="pictures/chDynamics/analysis6_4f_dissociation.pdf",
   headline="The model tells a clinician which kind of feature is atypical — local timing, global routing, or their coupling.",
   say="This is what I mean by a biological prism. Instead of one verdict, the model separates a patient's signal into where the abnormality lives — is it local timing, global network routing, or the coupling between them? That's a profile a clinician can reason about and act on. It's the difference between a number and an explanation."),

 dict(kind="bullets", title="What's next", bsize=21,
   bullets=["Single-trial decoding — use the full spiking representation, not trial averages",
            "On-chip on Loihi 2 — finally measure the energy I've only argued for",
            "A trained linear readout, and a formal bin-count study",
            "Clinical replication with the SHAPE lab on the layer-specific findings"],
   say="Where this goes next — and I'll be honest about what's still open. I worked with trial-averaged data, so single-trial decoding is the natural step. The energy argument has been exactly that, an argument, so the real test is on-chip, on neuromorphic hardware like Loihi, where I can finally measure it instead of motivating it. And the layer-specific clinical findings move into proper replication with the SHAPE lab. None of these is an afterthought — the architecture was built so each one is a direct extension rather than a redesign."),

 dict(kind="figure", title="The question, answered", img="pictures/conclusion/fig_defense_summary.pdf",
   headline="Stable. Compressible. Coherent. Interpretable. The architecture is the explanation.",
   say="So, back to the question I opened with. Can we recover hidden dynamics from a noisy trace and keep the explanation? The reservoir is stable — I measured it. The signal is compressible — six bins and sixty-four dimensions. The system is coherent — the coupling is real. And it's interpretable at four levels at once. With this design, the architecture isn't wrapped in an explanation. The architecture is the explanation."),

 dict(kind="close", title="Thank you.",
   sub="Andrew Lane  ·  with gratitude to Prof. Tang, Prof. Nelson and the SHAPE Lab, my committee, and my family.\nI'm happy to take your questions.",
   say="I'll close where I started — with that doctor and that patient. An instrument they can actually read changes what's possible in that room. Thank you to my advisor, Professor Tang; to Professor Nelson and the SHAPE lab; to my committee; and to my family. I'd be glad to take your questions."),

 # ============================ BACKUP ========================================
 dict(kind="section", kicker="BACKUP", title="Backup — for discussion"),

 dict(kind="bullets", title="Why a fixed reservoir, not a trained deep SNN", bsize=21,
   bullets=["Deep SNNs need backprop-through-time with surrogate gradients",
            "Forward/backward mismatch destabilizes long temporal horizons",
            "Universality: a fixed pool + linear readout already suffices",
            "Training the recurrence would erase the very transparency that is the point"],
   say="If you ask why I didn't train the reservoir: training a deep spiking network over hundreds of time steps needs backpropagation-through-time with surrogate gradients, and the forward and backward passes diverge over long horizons — it's brittle. A fixed reservoir sidesteps that entirely. More importantly, universality already guarantees a fixed pool plus a linear readout is enough, so I'm not giving anything up — and training the recurrence would destroy the very transparency that is the whole point."),

 dict(kind="figure", title="Spectral radius vs. measured stability",
   img="pictures/defense_audit/analysisA_1e_autonomous_vs_driven.pdf",
   headline="ρ(W) = 0.265 is a property of the weights; λ₁ = −0.054 is the system in operation.",
   say="If you press on stability: the autonomous spectral radius is zero point two six five, but that's the weights in isolation. What governs behavior is the driven exponent under the real input, and that's what I measured at minus zero point zero five four. I'm reporting the system in operation, not a textbook proxy."),

 dict(kind="figure", title="Memory-capacity regime and the choice of leak",
   img="pictures/defense_audit/analysisC_3d_memory_capacity_peak.pdf",
   headline="β = 0.05 sits at 91% of peak memory capacity, matched to the ERP window.",
   say="On why this leak and why six bins: the leak of zero point zero five sits on the memory-capacity plateau, at ninety-one percent of the peak, and it's set by matching the membrane timescale to the length of the ERP. Finer binning doesn't buy accuracy. It's a principled operating point, and a formal bin-count study is the defined next step."),

 dict(kind="figure", title="Embedding dimension, measured",
   img="pictures/defense_audit/analysisB_2e_takens_dimension.pdf",
   headline="False-nearest-neighbors puts the trajectory well under 64 dimensions.",
   say="Takens motivates the geometry, but I don't lean on it as a guarantee — I measured the embedding dimension with false-nearest-neighbors, and it sits well below the sixty-four I project into. The reservoir has comfortable margin."),

 dict(kind="figure", title="Full model comparison (N = 211)",
   img="pictures/chGraphNeuralNetworks/fig_05_gnn_comparison_211.pdf",
   headline="Same subject-disjoint protocol across every model.",
   say="Here's the complete two-hundred-eleven-subject comparison under one protocol. The reservoir matches the trained recurrent models with no trainable recurrence, and the ordering is exactly what the data-processing inequality predicts."),

 dict(kind="figure", title="Confusion structure",
   img="pictures/chGraphNeuralNetworks/confusion_matrix_best.pdf",
   headline="Errors concentrate between the two high-arousal extremes.",
   say="On where the errors fall: the confusions cluster between the two high-arousal classes, which points to arousal — not the sign of valence — carrying the separable axis. That's consistent with the pairwise analysis."),

 dict(kind="figure", title="Per-trial permutation null",
   img="pictures/defense_audit/analysisD_4d_permutation_nulls.pdf",
   headline="A separate channel permutation per trial — the strictest spatial null.",
   say="The strictest spatial control I ran: an independent channel permutation per trial, not one global shuffle. The classification still clears it. A flat classifier can't compensate for that."),

 dict(kind="figure", title="Four documented pivots",
   img="pictures/defense_audit/figG_failure_gallery.pdf",
   headline="Each negative result redirected the architecture.",
   say="And since methodology came up — these are the points where a result went against my hypothesis and I let it redirect the design rather than forcing the model through. The graph result is the clearest: I expected message passing to help, the theory said otherwise here, the data agreed, and that's what produced the design rule."),

 dict(kind="figure", title="Coupling detail",
   img="pictures/chSynthesis/fig7_A1_mean_coupling_matrix.pdf",
   headline="Structure–function coupling, the full matrix.",
   say="The full coupling matrix behind the kappa number, if you'd like to see how the dynamical and topological families align electrode by electrode."),

 dict(kind="figure", title="HC vs. MDD contrast",
   img="pictures/chDynamics/fig06_hc_vs_mdd.pdf",
   headline="Descriptor-level contrast, healthy controls vs. depression.",
   say="And the descriptor-level contrast between healthy controls and the depression group, for the clinical questions."),

 dict(kind="figure", title="ERP-window dependence",
   img="pictures/chDynamics/erp_window_comparison.pdf",
   headline="Discriminability by analysis window.",
   say="If you want the window analysis: this is how discriminability depends on which part of the epoch we read, and it lines up with the early components rather than the late drift the black box leaned on."),

 dict(kind="figure", title="Layers are complementary, not redundant",
   img="pictures/chDynamics/analysis6_3_four_path_comparison.pdf",
   headline="Combining families never beats the better one — but each disorder loads on a different family.",
   say="On whether the layers are just redundant: combining the temporal and spatial families never beats the better one alone, yet different disorders load on different families. That's the definition of complementary information, and it's why the decomposition is worth keeping."),

 dict(kind="figure", title="Coupling by affective category",
   img="pictures/chSynthesis/fig7_A4_kappa_by_category.pdf",
   headline="κ structure across stimulus subcategories.",
   say="The coupling broken out by affective subcategory, if you want to see how the structure-function relationship shifts with the stimulus."),

 dict(kind="figure", title="Layer-specific clinical sensitivity",
   img="pictures/chDynamics/fig03_clinical_heatmap.pdf",
   headline="Which descriptor responds to which diagnosis.",
   say="And the full heatmap behind the disorder-layer table — which descriptor moves for which diagnosis, with substance use standing out, as the earlier slide showed."),

 dict(kind="table", title="Three generations of neural computation", img="tbl_generations",
   say="If the framing comes up: first-generation units are threshold logic with no notion of time; second-generation deep nets abstract time away into firing rates; and the third generation — spiking — puts time back as the carrier. That's the family ARSPI-Net belongs to, and why it fits a signal whose information is in its timing."),
]

def main():
    prs = Presentation()
    prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
    n = 0
    for sl in SLIDES:
        if sl["kind"] not in ("title", "section", "close"):
            n += 1
        render(prs, sl, n)
    os.makedirs(os.path.dirname(OUT_PPTX), exist_ok=True)
    prs.save(OUT_PPTX)

    # transcript
    with open(OUT_TXT, "w") as f:
        f.write("# ARSPI-Net — Doctoral Defense: Spoken Transcript\n\n")
        f.write("*Andrew Lane · Stony Brook ECE · ~40 minutes. "
                "Word-for-word script; the same text is in each slide's speaker notes.*\n\n---\n\n")
        idx = 0
        for sl in SLIDES:
            if sl["kind"] == "section":
                f.write(f"\n## ◆ {sl['title']}\n\n")
                if sl.get("say"):
                    f.write(f"*{sl['say']}*\n\n")
                continue
            if sl["kind"] not in ("title", "close"):
                idx += 1
                head = f"### Slide {idx} — {sl.get('title') or sl.get('kicker','')}"
            elif sl["kind"] == "title":
                head = "### Title"
            else:
                head = "### Closing"
            f.write(f"{head}\n\n{sl.get('say','')}\n\n")

    nmain = nback = 0; seen_backup = False
    for s in SLIDES:
        if s.get("kicker") == "BACKUP":
            seen_backup = True
        if s["kind"] in ("title", "section", "close"):
            continue
        if seen_backup: nback += 1
        else: nmain += 1
    total = len(SLIDES)
    print(f"Saved {OUT_PPTX}")
    print(f"Saved {OUT_TXT}")
    print(f"Slides: {total} total  |  ~{nmain} numbered main  |  ~{nback} backup")
    if missing:
        print("MISSING ASSETS:", sorted(set(missing)))
    else:
        print("All assets resolved.")

if __name__ == "__main__":
    main()
